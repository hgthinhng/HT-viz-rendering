import re
from html.parser import HTMLParser
from pptx import Presentation

class TextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self.text = []
        self.skip = False
    def handle_starttag(self, tag, attrs):
        if tag in ('style', 'script'):
            self.skip = True
    def handle_endtag(self, tag):
        if tag in ('style', 'script'):
            self.skip = False
    def handle_data(self, data):
        if not self.skip:
            t = data.strip()
            if t:
                self.text.append(t)

def extract_text(path):
    with open(path, encoding='utf-8') as f:
        html = f.read()
    p = TextExtractor()
    p.feed(html)
    return ' '.join(p.text)

def nowhite(s):
    return re.sub(r'\s+', '', s)

src_text = extract_text('slide-test-nosvg.html')
src_chars = nowhite(src_text)

prs = Presentation('output.pptx')
out_parts = []
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            out_parts.append(shape.text_frame.text)
out_text = ' '.join(out_parts)
out_chars = nowhite(out_text)

print(f"Source visible text chars (no-svg slide, whitespace stripped): {len(src_chars)}")
print(f"Output PPTX text chars (whitespace stripped):                  {len(out_chars)}")
print(f"Ratio: {len(out_chars)}/{len(src_chars)} = {len(out_chars)/len(src_chars)*100:.1f}%")
print()

# What's missing = table content
missing = src_chars
for ch_run in re.findall(r'.', out_chars):
    pass
# simple diff: characters in src not accounted for (crude but illustrative)
import difflib
sm = difflib.SequenceMatcher(None, src_chars, out_chars)
missing_spans = []
for tag, i1, i2, j1, j2 in sm.get_opcodes():
    if tag in ('delete', 'replace'):
        missing_spans.append(src_chars[i1:i2])
print("Missing (present in source HTML, absent from PPTX):")
print(''.join(missing_spans))
