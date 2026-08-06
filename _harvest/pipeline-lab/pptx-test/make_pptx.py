from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# editable title text box
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
tf = title_box.text_frame
tf.text = "Bao cao tai chinh quy 2 - Test native PPTX"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0x1e, 0x29, 0x3b)

# editable body text box with Vietnamese diacritics
body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(1.5))
btf = body_box.text_frame
btf.word_wrap = True
p = btf.paragraphs[0]
p.text = "Chu co dau: ừ ộ ẫ ợ ữ ể ỗ - kiem tra font tren PowerPoint that."
p.font.size = Pt(16)

# native chart (real chart object, editable data, not an image)
chart_data = CategoryChartData()
chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
chart_data.add_series('Doanh thu (ty VND)', (120, 145, 160, 190))
chart_data.add_series('Loi nhuan (ty VND)', (20, 28, 31, 40))

x, y, cx, cy = Inches(0.5), Inches(3.2), Inches(8), Inches(4)
graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
chart = graphic_frame.chart
chart.has_legend = True

prs.save('native_test.pptx')
print("Saved native_test.pptx")
